"""
Turn raw file bytes into plain text. This is the piece Iron Bot is missing
today -- it currently only sees an attachment's *filename*, never its contents.

Supported: PDF (text layer + OCR fallback for scans), Word (.docx), Excel
(.xlsx), PowerPoint (.pptx), images (OCR), and plain text / csv / md.

Every extractor degrades gracefully: if an optional dependency or the OCR
binary is missing, it returns "" instead of raising, and logs why.
"""
import io
import logging
import os

logger = logging.getLogger("rag.extract")

# Extensions we will attempt to extract. Anything else is skipped.
TEXT_EXTS = {".txt", ".md", ".markdown", ".csv", ".tsv", ".log", ".json"}
PDF_EXTS = {".pdf"}
DOCX_EXTS = {".docx"}
XLSX_EXTS = {".xlsx", ".xlsm"}
PPTX_EXTS = {".pptx"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"}

SUPPORTED_EXTS = (
    TEXT_EXTS | PDF_EXTS | DOCX_EXTS | XLSX_EXTS | PPTX_EXTS | IMAGE_EXTS
)


def is_supported(filename):
    return os.path.splitext(filename or "")[1].lower() in SUPPORTED_EXTS


def extract_text(file_bytes, filename):
    """Dispatch on extension. Returns extracted text (possibly "")."""
    ext = os.path.splitext(filename or "")[1].lower()
    try:
        if ext in TEXT_EXTS:
            return _decode_text(file_bytes)
        if ext in PDF_EXTS:
            return _extract_pdf(file_bytes)
        if ext in DOCX_EXTS:
            return _extract_docx(file_bytes)
        if ext in XLSX_EXTS:
            return _extract_xlsx(file_bytes)
        if ext in PPTX_EXTS:
            return _extract_pptx(file_bytes)
        if ext in IMAGE_EXTS:
            return _ocr_image(file_bytes)
    except Exception as e:  # never let one bad file kill the whole ingest
        logger.warning("extract failed for %s: %s", filename, str(e)[:200])
        return ""
    logger.info("unsupported file type skipped: %s", filename)
    return ""


def _decode_text(file_bytes):
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return file_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")


def _extract_pdf(file_bytes):
    """Text layer via PyMuPDF (already a repo dependency). OCR pages that
    come back empty (scanned PDFs)."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.warning("PyMuPDF not installed; cannot read PDFs")
        return ""
    out = []
    needs_ocr = []
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for i, page in enumerate(doc):
            txt = page.get_text("text").strip()
            if txt:
                out.append(txt)
            else:
                needs_ocr.append(i)
        if needs_ocr:
            for i in needs_ocr:
                ocr = _ocr_pdf_page(doc[i])
                if ocr:
                    out.append(ocr)
    return "\n\n".join(out).strip()


def _ocr_pdf_page(page):
    try:
        import fitz  # noqa: F401
        pix = page.get_pixmap(matrix=__import__("fitz").Matrix(2, 2))
        return _ocr_image(pix.tobytes("png"))
    except Exception as e:
        logger.warning("pdf page OCR failed: %s", str(e)[:120])
        return ""


def _ocr_image(file_bytes):
    """OCR via pytesseract. Requires the `tesseract` binary (installed in CI)."""
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        logger.info("pytesseract/Pillow not installed; skipping image OCR")
        return ""
    try:
        img = Image.open(io.BytesIO(file_bytes))
        return pytesseract.image_to_string(img).strip()
    except Exception as e:
        logger.warning("image OCR failed: %s", str(e)[:120])
        return ""


def _extract_docx(file_bytes):
    try:
        import docx  # python-docx
    except ImportError:
        logger.warning("python-docx not installed; cannot read .docx")
        return ""
    d = docx.Document(io.BytesIO(file_bytes))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _extract_xlsx(file_bytes):
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl not installed; cannot read .xlsx")
        return ""
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _extract_pptx(file_bytes):
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed; cannot read .pptx")
        return ""
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
    return "\n".join(parts).strip()
